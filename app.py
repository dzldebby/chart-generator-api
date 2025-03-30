from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import os
import tempfile
import pandas as pd
import io
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.util import Inches
import base64
from datetime import datetime

app = Flask(__name__)
# Enable CORS for all routes with all origins
CORS(app, resources={r"/*": {"origins": "*"}})

# Store generated files temporarily (in a real app, use a proper storage solution)
GENERATED_FILES = {}

# Chart type mapping
CHART_TYPE_MAP = {
    'pie': XL_CHART_TYPE.PIE,
    'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
    'line': XL_CHART_TYPE.LINE,
    'scatter': XL_CHART_TYPE.XY_SCATTER,
    'area': XL_CHART_TYPE.AREA
}

@app.route('/', methods=['GET'])
def home():
    return jsonify({
        "status": "online",
        "message": "Chart generation service is running"
    })

@app.route('/generate-chart', methods=['POST'])
def generate_chart():
    try:
        # Check if files are in the request
        if 'dataFile' not in request.files:
            return jsonify({"error": "No data file provided"}), 400

        # Get form fields
        chart_type = request.form.get('chartType', 'bar')
        slide_position = int(request.form.get('slidePosition', 1))
        
        # Get the data file
        data_file = request.files['dataFile']
        
        # Get template file if provided
        template_file = None
        if 'templateFile' in request.files and request.files['templateFile'].filename:
            template_file = request.files['templateFile']
        
        # Create a temp directory for file operations
        with tempfile.TemporaryDirectory() as temp_dir:
            # Save the data file
            data_file_path = os.path.join(temp_dir, data_file.filename)
            data_file.save(data_file_path)
            
            # Process the data file
            if data_file_path.endswith('.csv'):
                df = pd.read_csv(data_file_path)
            elif data_file_path.endswith(('.xlsx', '.xls')):
                df = pd.read_excel(data_file_path)
            else:
                return jsonify({"error": f"Unsupported file format: {data_file.filename}"}), 400
                
            # Create or load presentation
            if template_file:
                template_path = os.path.join(temp_dir, template_file.filename)
                template_file.save(template_path)
                prs = Presentation(template_path)
            else:
                prs = Presentation()
            
            # Make sure we have enough slides
            while len(prs.slides) < slide_position:
                # Try to find a blank layout with title
                layout_idx = 5  # Default to layout 5 (usually blank with title)
                for idx, layout in enumerate(prs.slide_layouts):
                    if layout.name and ('blank' in layout.name.lower() or 'title' in layout.name.lower()):
                        layout_idx = idx
                        break
                prs.slides.add_slide(prs.slide_layouts[layout_idx])
                
            # Get the target slide (0-based index)
            slide_idx = slide_position - 1
            slide = prs.slides[slide_idx]
            
            # Create chart data
            chart_data = ChartData()
            
            # Get first column as categories and second as values
            # Handle different data structures
            if df.shape[1] < 2:
                return jsonify({"error": "Data must have at least two columns (categories and values)"}), 400
                
            categories = df.iloc[:, 0].tolist()
            
            # Handle percentage values if they're strings with '%'
            values = []
            for val in df.iloc[:, 1]:
                if isinstance(val, str) and '%' in val:
                    try:
                        values.append(float(val.strip('%')))
                    except ValueError:
                        values.append(0)
                else:
                    try:
                        values.append(float(val))
                    except (ValueError, TypeError):
                        values.append(0)
                
            chart_data.categories = categories
            chart_data.add_series('Values', values)
            
            # Add chart to slide
            x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4)
            chart = slide.shapes.add_chart(
                CHART_TYPE_MAP.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED),
                x, y, cx, cy, chart_data
            ).chart
            
            # Customize chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            
            # Add data labels for pie charts
            if chart_type == 'pie':
                chart.plots[0].has_data_labels = True
                chart.plots[0].data_labels.number_format = '0.0"%"'
            
            # Save presentation to memory
            pptx_buffer = io.BytesIO()
            prs.save(pptx_buffer)
            pptx_buffer.seek(0)
            
            # Generate a unique ID for this file
            file_id = f"{datetime.now().strftime('%Y%m%d%H%M%S')}_{chart_type}"
            filename = f"chart_{file_id}.pptx"
            
            # Store the buffer in our dict (in a real app, use proper storage)
            GENERATED_FILES[file_id] = {
                'buffer': pptx_buffer,
                'filename': filename,
                'created': datetime.now(),
                'content_type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            }
            
            # For a real preview, you'd need to convert the first slide to an image
            # For simplicity, we'll just use a placeholder
            preview_url = f"https://via.placeholder.com/800x450.png?text=Chart+Preview+{chart_type}"
            
            # Return URLs for preview and download
            return jsonify({
                "message": "Chart generated successfully",
                "downloadUrl": f"/download-chart/{file_id}",
                "previewUrl": preview_url,
                "chartType": chart_type
            })
            
    except Exception as e:
        import traceback
        print(traceback.format_exc())
        return jsonify({"error": str(e)}), 500

@app.route('/download-chart/<file_id>', methods=['GET'])
def download_chart(file_id):
    if file_id not in GENERATED_FILES:
        return jsonify({"error": "File not found"}), 404
        
    file_data = GENERATED_FILES[file_id]
    file_data['buffer'].seek(0)
    
    return send_file(
        file_data['buffer'],
        mimetype=file_data['content_type'],
        as_attachment=True,
        download_name=file_data['filename']
    )

# Clean up old files periodically (in a real app, this would be a background task)
@app.route('/cleanup', methods=['GET'])
def cleanup():
    now = datetime.now()
    removed = 0
    
    for file_id in list(GENERATED_FILES.keys()):
        file_data = GENERATED_FILES[file_id]
        age = (now - file_data['created']).total_seconds()
        
        # Remove files older than 1 hour
        if age > 3600:
            del GENERATED_FILES[file_id]
            removed += 1
    
    return jsonify({
        "message": f"Cleanup complete. Removed {removed} old files.",
        "remaining": len(GENERATED_FILES)
    })

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port)